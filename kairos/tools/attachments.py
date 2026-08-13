"""Attachment extraction and explicit note-creation workflow."""

from __future__ import annotations

import json
import secrets
import sqlite3
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path

import requests

from kairos.channels.feishu import tenant_token, user_feishu_request
from kairos.infrastructure.settings import DB_PATH
from kairos.infrastructure.llm import build_client, model_name
from kairos.tools.text import plain_text

http = requests.Session()
http.headers["User-Agent"] = "Kairós/1.0"
llm = build_client()


def message_resource(message_id: str, file_key: str) -> bytes:
    """Download an attachment sent in a Feishu message."""
    response = http.get(
        f"https://open.feishu.cn/open-apis/im/v1/messages/{message_id}/resources/{file_key}",
        headers={"Authorization": f"Bearer {tenant_token()}"},
        params={"type": "file"},
        timeout=90,
    )
    response.raise_for_status()
    return response.content


def extract_attachment(name: str, blob: bytes) -> str:
    """Extract plain text from the supported common office formats."""
    suffix = Path(name).suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return blob.decode("utf-8", errors="replace")
    if suffix == ".pdf":
        from pypdf import PdfReader

        return "\n".join((page.extract_text() or "") for page in PdfReader(BytesIO(blob)).pages)
    if suffix == ".docx":
        from docx import Document

        return "\n".join(paragraph.text for paragraph in Document(BytesIO(blob)).paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation

        return "\n".join(
            shape.text
            for slide in Presentation(BytesIO(blob)).slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
    raise ValueError("暂只支持 PDF、Word、PPT 和纯文本附件。")


def prepare_note(message_id: str, content: str) -> str:
    """Create a preview; writing a cloud document always requires confirmation."""
    info = json.loads(content)
    name, key = info.get("file_name", "附件"), info.get("file_key", "")
    if not key:
        return "老师，附件缺少 file_key，无法下载。"
    text = extract_attachment(name, message_resource(message_id, key))[:60_000]
    if not text.strip():
        return "老师，附件中未提取到可读文本。"

    prompt = "用中文纯文本把下面附件整理为研究笔记：主题、核心要点、方法或证据、待办。不要 Markdown。\n\n" + text
    note = plain_text(
        llm.chat.completions.create(model=model_name(), messages=[{"role": "user", "content": prompt}], temperature=0.2).choices[0].message.content
        or ""
    )
    code = secrets.token_hex(3).upper()
    with sqlite3.connect(DB_PATH) as connection:
        connection.execute(
            "INSERT OR REPLACE INTO pending_notes VALUES (?, ?, ?, ?)",
            (code, f"笔记：{Path(name).stem}", note, datetime.now(UTC).isoformat()),
        )
    return f"老师，已完成《{name}》的笔记预览：\n\n{note[:5000]}\n\n如需写入飞书云文档，请回复：确认笔记 {code}"


def create_note(code: str) -> str:
    """Write a previously previewed note to Feishu after explicit confirmation."""
    with sqlite3.connect(DB_PATH) as connection:
        row = connection.execute("SELECT title, content FROM pending_notes WHERE code=?", (code.upper(),)).fetchone()
    if not row:
        return "老师，没有找到待确认的笔记编号。"

    title, content = row
    created = user_feishu_request("POST", "/docx/v1/documents", json={"title": title}).get("data", {}).get("document", {})
    document_id = created.get("document_id")
    if not document_id:
        return "老师，创建云文档失败。"
    children = [
        {"block_type": 2, "text": {"elements": [{"text_run": {"content": content[index:index + 1400], "text_element_style": {}}}]}}
        for index in range(0, len(content), 1400)
    ][:40]
    user_feishu_request(
        "POST",
        f"/docx/v1/documents/{document_id}/blocks/{document_id}/children",
        json={"children": children, "index": -1},
    )
    with sqlite3.connect(DB_PATH) as connection:
        connection.execute("DELETE FROM pending_notes WHERE code=?", (code.upper(),))
    return f"老师，笔记已创建：https://my.feishu.cn/docx/{document_id}"
